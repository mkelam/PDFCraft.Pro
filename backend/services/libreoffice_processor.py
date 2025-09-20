"""
LibreOffice PDF to PPT Conversion Service
High-performance PDF to PowerPoint conversion using LibreOffice headless mode
"""

import asyncio
import subprocess
import tempfile
import shutil
import os
from pathlib import Path
from typing import Optional, Dict, Any
import time
import uuid

class LibreOfficeProcessor:
    """
    LibreOffice headless processor for PDF-to-PPT conversion
    Targets sub-6 second conversion performance
    """

    def __init__(self):
        self.temp_dir = Path(tempfile.gettempdir()) / "libreoffice_processor"
        self.temp_dir.mkdir(exist_ok=True)
        self.libreoffice_path = self._find_libreoffice()
        self.max_concurrent_conversions = 3  # Optimal for performance
        self.conversion_semaphore = asyncio.Semaphore(self.max_concurrent_conversions)

    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice executable on the system"""
        possible_paths = [
            # Windows paths
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
            # Linux paths
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",
            "/snap/bin/libreoffice",
            # macOS paths
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            # Docker/Container path
            "/opt/libreoffice/program/soffice"
        ]

        # Check environment variable first
        env_path = os.getenv("LIBREOFFICE_PATH")
        if env_path and Path(env_path).exists():
            return env_path

        # Check common installation paths
        for path in possible_paths:
            if Path(path).exists():
                return path

        # Try to find in PATH on Windows
        try:
            result = subprocess.run(["where", "soffice.exe"],
                                  capture_output=True, text=True, timeout=5)
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip().split('\n')[0]  # Take first result
        except:
            pass

        # Try to find libreoffice
        try:
            result = subprocess.run(["where", "libreoffice"],
                                  capture_output=True, text=True, timeout=5)
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip().split('\n')[0]  # Take first result
        except:
            pass

        return None

    async def initialize(self) -> bool:
        """Initialize LibreOffice processor and verify installation"""
        if not self.libreoffice_path:
            print("LibreOffice not found on system")
            print("Install LibreOffice or set LIBREOFFICE_PATH environment variable")
            return False

        # LibreOffice executable found - assume it works
        # Skip version check as it can hang on Windows
        print(f"LibreOffice found at: {self.libreoffice_path}")
        print(f"Temp directory: {self.temp_dir}")
        print("Skipping version check - will test during first conversion")
        return True

    async def convert_pdf_to_ppt(self, pdf_path: str, job_id: Optional[str] = None) -> Dict[str, Any]:
        """
        Convert PDF to PowerPoint presentation

        Args:
            pdf_path: Path to source PDF file
            job_id: Optional job ID for tracking

        Returns:
            Dictionary with conversion results
        """
        async with self.conversion_semaphore:
            start_time = time.time()
            job_id = job_id or str(uuid.uuid4())

            try:
                # Validate input file
                pdf_file = Path(pdf_path)
                if not pdf_file.exists():
                    raise FileNotFoundError(f"PDF file not found: {pdf_path}")

                if not pdf_file.suffix.lower() == '.pdf':
                    raise ValueError(f"Input file must be PDF, got: {pdf_file.suffix}")

                # Create job-specific temp directory
                job_temp_dir = self.temp_dir / f"job_{job_id}"
                job_temp_dir.mkdir(exist_ok=True)

                # Copy PDF to temp directory (LibreOffice works better with local files)
                temp_pdf = job_temp_dir / f"input_{job_id}.pdf"
                shutil.copy2(pdf_path, temp_pdf)

                print(f"Starting PDF to PPT conversion: {pdf_file.name}")

                # Run LibreOffice conversion
                try:
                    result = await self._convert_with_libreoffice(
                        input_file=temp_pdf,
                        output_dir=job_temp_dir,
                        job_id=job_id
                    )
                except RuntimeError as e:
                    if "no export filter" in str(e) or "platform independent libraries" in str(e):
                        # LibreOffice doesn't support PDF-to-PPT conversion
                        # Use alternative conversion method
                        print("LibreOffice PDF to PPT not supported - using alternative conversion method")
                        try:
                            result = await self._convert_pdf_to_pptx_alternative(
                                input_file=temp_pdf,
                                output_dir=job_temp_dir,
                                job_id=job_id
                            )
                        except Exception as alt_error:
                            print(f"Alternative conversion also failed: {alt_error}")
                            raise RuntimeError(f"PDF to PowerPoint conversion failed. LibreOffice reports: {str(e)}. Alternative method failed: {str(alt_error)}")
                    else:
                        raise

                conversion_time = time.time() - start_time

                # Find generated PPT file
                ppt_files = list(job_temp_dir.glob("*.ppt*"))
                if not ppt_files:
                    # No output file generated - this is a real failure
                    raise RuntimeError("LibreOffice conversion completed but no PPT file was generated. Conversion failed.")

                output_file = ppt_files[0]
                final_output = self.temp_dir / f"converted_{job_id}.pptx"

                # Move to final location
                shutil.move(str(output_file), str(final_output))

                # Get file sizes for metrics
                input_size = pdf_file.stat().st_size
                output_size = final_output.stat().st_size

                # Cleanup job temp directory
                shutil.rmtree(job_temp_dir, ignore_errors=True)

                # Performance check
                target_time = 6.0  # 6 second target
                is_within_target = conversion_time <= target_time

                print(f"{'SUCCESS' if is_within_target else 'SLOW'} Conversion completed:")
                print(f"  Time: {conversion_time:.2f}s (Target: <6s)")
                print(f"  Input: {input_size:,} bytes ({pdf_file.name})")
                print(f"  Output: {output_size:,} bytes ({final_output.name})")
                print(f"  Performance: {'WITHIN TARGET' if is_within_target else 'EXCEEDED TARGET'}")

                return {
                    "success": True,
                    "job_id": job_id,
                    "input_file": str(pdf_path),
                    "output_file": str(final_output),
                    "conversion_time_seconds": conversion_time,
                    "input_size_bytes": input_size,
                    "output_size_bytes": output_size,
                    "within_performance_target": is_within_target,
                    "libreoffice_version": await self._get_libreoffice_version()
                }

            except Exception as e:
                conversion_time = time.time() - start_time
                error_msg = f"PDF to PPT conversion failed: {str(e)}"

                print(f"ERROR: {error_msg}")
                print(f"  Failed after: {conversion_time:.2f}s")

                # Cleanup on error
                job_temp_dir = self.temp_dir / f"job_{job_id}"
                if job_temp_dir.exists():
                    shutil.rmtree(job_temp_dir, ignore_errors=True)

                return {
                    "success": False,
                    "job_id": job_id,
                    "error": error_msg,
                    "conversion_time_seconds": conversion_time
                }

    async def _convert_with_libreoffice(self, input_file: Path, output_dir: Path, job_id: str) -> subprocess.CompletedProcess:
        """
        Execute LibreOffice conversion command
        """
        # LibreOffice command for PDF-to-PPT conversion
        command = [
            "--headless",                    # Run without GUI
            "--convert-to", "pptx",         # Convert to PowerPoint format
            "--outdir", str(output_dir),    # Output directory
            str(input_file)                 # Input PDF file
        ]

        print(f"LibreOffice command: {' '.join(command[:4])} [...]")

        # Execute with timeout (should complete in <6 seconds)
        result = await self._run_libreoffice_command(command, timeout=30)

        # Check for LibreOffice errors even when exit code is 0
        stderr_output = result.stderr.decode() if result.stderr else ""

        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed (exit code {result.returncode}): {stderr_output}")

        # Check for "no export filter" error even with successful exit code
        if "no export filter" in stderr_output or "platform independent libraries" in stderr_output:
            raise RuntimeError(f"LibreOffice PDF-to-PPT conversion not supported: {stderr_output}")

        return result

    async def _run_libreoffice_command(self, command: list, timeout: int = 30) -> subprocess.CompletedProcess:
        """
        Run LibreOffice command using thread pool executor (fixes asyncio subprocess issues)
        Enhanced with diagnostic logging and Windows-specific fixes
        """
        try:
            # Prepend the LibreOffice executable path to the command
            full_command = [self.libreoffice_path] + command

            # Enhanced diagnostic logging
            print(f"DIAGNOSTIC: Full command: {' '.join(full_command)}")
            print(f"DIAGNOSTIC: Working directory: {os.getcwd()}")
            print(f"DIAGNOSTIC: LibreOffice path exists: {Path(self.libreoffice_path).exists()}")
            print(f"DIAGNOSTIC: Current user: {os.getenv('USERNAME', 'Unknown')}")
            print(f"DIAGNOSTIC: Environment PATH: {os.getenv('PATH', 'Not found')[:200]}...")

            # Use thread pool executor to run synchronous subprocess (avoids asyncio subprocess issues)
            print("DIAGNOSTIC: Using thread pool executor for subprocess execution")

            def run_sync_subprocess():
                """Run subprocess synchronously in thread pool"""
                import subprocess
                return subprocess.run(
                    full_command,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=timeout,
                    env=os.environ.copy()
                )

            # Run in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(None, run_sync_subprocess)

            # Enhanced diagnostic output
            print(f"DIAGNOSTIC: Process return code: {result.returncode}")
            print(f"DIAGNOSTIC: Stdout length: {len(result.stdout) if result.stdout else 0} bytes")
            print(f"DIAGNOSTIC: Stderr length: {len(result.stderr) if result.stderr else 0} bytes")

            if result.stdout:
                stdout_text = result.stdout.decode('utf-8', errors='ignore')
                print(f"DIAGNOSTIC: Stdout content: {stdout_text[:500]}...")

            if result.stderr:
                stderr_text = result.stderr.decode('utf-8', errors='ignore')
                print(f"DIAGNOSTIC: Stderr content: {stderr_text[:500]}...")

            return result

        except (asyncio.TimeoutError, subprocess.TimeoutExpired) as e:
            # Handle timeout from either asyncio or subprocess
            print(f"DIAGNOSTIC: Process timed out after {timeout}s: {e}")
            raise RuntimeError(f"LibreOffice command timed out after {timeout}s")

        except Exception as e:
            print(f"DIAGNOSTIC: Subprocess creation/execution error: {e}")
            print(f"DIAGNOSTIC: Error type: {type(e).__name__}")
            print(f"DIAGNOSTIC: Command that failed: {' '.join(full_command)}")
            raise RuntimeError(f"LibreOffice subprocess execution failed: {e}")

    async def _get_libreoffice_version(self) -> str:
        """Get LibreOffice version for metadata"""
        try:
            result = await self._run_libreoffice_command(["--version"], timeout=5)
            if result.returncode == 0:
                return result.stdout.decode().strip()
            return "Unknown"
        except:
            return "Unknown"

    async def _convert_pdf_to_pptx_alternative(self, input_file: Path, output_dir: Path, job_id: str):
        """
        Alternative PDF to PPTX conversion using pdf2image + python-pptx
        Creates real PowerPoint presentations with PDF pages as images
        """
        try:
            print(f"Starting alternative PDF to PPTX conversion for {input_file}")

            # Check if required packages are available
            try:
                from pdf2image import convert_from_path
                from pptx import Presentation
                from pptx.util import Inches
                import io
            except ImportError as e:
                print(f"Required packages not installed: {e}")
                print("Please install: pip install pdf2image python-pptx Pillow")
                raise RuntimeError(f"Alternative conversion dependencies missing: {e}")

            # Convert PDF pages to images
            print("Converting PDF pages to images...")
            try:
                images = convert_from_path(str(input_file), dpi=150, fmt='PNG')
                print(f"Successfully converted {len(images)} pages to images")
            except Exception as e:
                print(f"Failed to convert PDF to images: {e}")
                raise RuntimeError(f"PDF to image conversion failed: {e}")

            # Create PowerPoint presentation
            print("Creating PowerPoint presentation...")
            prs = Presentation()

            # Set presentation dimensions (16:9 aspect ratio)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            for i, img in enumerate(images):
                print(f"Adding slide {i+1}/{len(images)}")

                # Add slide with blank layout
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)

                # Convert PIL image to bytes
                img_stream = io.BytesIO()
                img.save(img_stream, format='PNG')
                img_stream.seek(0)

                # Calculate image dimensions to fit slide while maintaining aspect ratio
                img_width, img_height = img.size
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # Calculate scaling factor
                width_ratio = slide_width / img_width
                height_ratio = slide_height / img_height
                scale_factor = min(width_ratio, height_ratio)

                # Calculate final dimensions and position
                final_width = int(img_width * scale_factor)
                final_height = int(img_height * scale_factor)
                left = (slide_width - final_width) // 2
                top = (slide_height - final_height) // 2

                # Add image to slide
                slide.shapes.add_picture(
                    img_stream,
                    left, top,
                    final_width, final_height
                )

            # Save the presentation
            output_file = output_dir / f"input_{job_id}.pptx"
            prs.save(str(output_file))

            print(f"Alternative conversion completed: {output_file}")
            print(f"Created presentation with {len(images)} slides")

            # Validate the created file
            if not self._validate_pptx_file(output_file):
                raise RuntimeError("Generated PPTX file failed validation")

            return output_file

        except Exception as e:
            print(f"Alternative PDF to PPTX conversion failed: {e}")
            raise

    def _validate_pptx_file(self, file_path: Path) -> bool:
        """Validate that the output file is a real PPTX file"""
        try:
            # Check file size (should be > 1KB for real PPTX)
            if file_path.stat().st_size < 1024:
                print(f"PPTX validation failed: file too small ({file_path.stat().st_size} bytes)")
                return False

            # Verify ZIP structure
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Real PPTX should contain these essential files
                required_files = ['[Content_Types].xml', 'ppt/presentation.xml']
                zip_contents = zip_file.namelist()

                for required_file in required_files:
                    if required_file not in zip_contents:
                        print(f"PPTX validation failed: missing {required_file}")
                        return False

            print(f"PPTX validation passed: {file_path}")
            return True

        except Exception as e:
            print(f"PPTX validation error: {e}")
            return False

    async def cleanup(self):
        """Cleanup temporary files and resources"""
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir, ignore_errors=True)
        print("LibreOffice processor cleanup completed")

# Global instance for use in FastAPI
libreoffice_processor = LibreOfficeProcessor()